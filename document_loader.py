# utils/document_loader.py

import os
import textract
from PyPDF2 import PdfReader
from docx import Document
import csv
import pptx

def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        reader = PdfReader(file_path)
        return "\n".join([page.extract_text() for page in reader.pages])

    elif ext == ".docx":
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    elif ext == ".txt" or ext == ".md":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    elif ext == ".csv":
        with open(file_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            return "\n".join([", ".join(row) for row in reader])

    elif ext == ".pptx":
        prs = pptx.Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    else:
        raise ValueError("Unsupported file format")
